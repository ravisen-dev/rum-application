from pptx import Presentation
from pptx.util import Inches, Pt

prs = Presentation()
prs.slide_width = Inches(13.33)
prs.slide_height = Inches(7.5)

def add_slide(title, bullets):
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = title
    body = slide.shapes.placeholders[1].text_frame
    body.clear()
    body.text = bullets[0]
    for item in bullets[1:]:
        p = body.add_paragraph()
        p.text = item
        p.level = 0
        p.font.size = Pt(24)
    return slide

slide = prs.slides.add_slide(prs.slide_layouts[0])
slide.shapes.title.text = 'Real User Monitoring (RUM) Solution'
slide.placeholders[1].text = 'Business overview, benefits, application flow, architecture, and technology stack.'

add_slide('Why RUM Matters', [
    'RUM captures real user behavior and performance in production.',
    'Provides visibility into actual user experience, not just backend metrics.',
    'Helps prioritize improvements based on real customer impact.',
    'Supports faster problem resolution with actionable insights.'
])

add_slide('Benefits over Traditional Observability', [
    'Real User Monitoring focuses on customer experience rather than system metrics.',
    'Traditional tools like Grafana visualize infrastructure metrics, but RUM shows front-end latency, page load, and errors in context.',
    'RUM identifies problems in the application path that observability dashboards may miss.',
    'Combines behavioral telemetry, performance, errors, and network activity in one view.'
])

add_slide('Applications Flow', [
    'User interacts with the frontend application.',
    'RUM SDK collects page views, network requests, errors, web vitals, and custom events.',
    'Telemetry is sent to the backend RUM API ingestion endpoint.',
    'Backend stores and processes telemetry data for analysis.',
    'Dashboard presents aggregated application insights to stakeholders.'
])

add_slide('Architecture Overview', [
    'Frontend apps use the RUM SDK to instrument user activity.',
    'RUM API backend ingests telemetry and stores it in PostgreSQL.',
    'Dashboard reads telemetry from the backend API for visualization.',
    'CORS-enabled local dev setup supports web apps and dashboard on separate ports.'
])

add_slide('Technology Stack', [
    'rum-api: ASP.NET Core Minimal API, EF Core, PostgreSQL.',
    'rum-dashboard: Angular 21 frontend dashboard.',
    'rum-sdk: TypeScript SDK for browser telemetry collection.',
    'rum-sdk-todos: Angular sample app demonstrating SDK integration.',
    'Dev tools: Node.js, npm, dotnet, TypeScript, Angular CLI.'
])

prs.save('RUM-Business-Deck.pptx')
print('Created RUM-Business-Deck.pptx')
